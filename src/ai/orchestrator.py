from typing import List, Dict, Any
import aiohttp
import asyncio
from langchain_openai import ChatOpenAI
from langchain.prompts import PromptTemplate
from langchain_core.runnables import RunnableSequence
import json
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
import matplotlib.pyplot as plt
import seaborn as sns

class AIOrchestrator:
    def __init__(self, api_key: str):
        self.llm = ChatOpenAI(
            temperature=0,
            model="gpt-4",  # Using GPT-4 for better reasoning
            openai_api_key=api_key
        )
        
        # Create research chain
        self.research_prompt = PromptTemplate.from_template(
            """Research the following company and provide key information:
            Company Domain: {domain}
            
            Focus on:
            1. Industry and business type
            2. Company size and locations
            3. Technologies they likely use
            4. Recent security incidents or breaches
            5. Regulatory requirements
            6. Critical assets and crown jewels
            
            Format the response as a JSON object with these keys:
            - company_info: Basic company information
            - tech_stack: Likely technology stack
            - security_history: Past incidents
            - compliance: Regulatory requirements
            - critical_assets: Key assets to protect
            - risk_factors: Industry-specific risks
            """
        )
        self.research_chain = self.research_prompt | self.llm
        
        # Create planning chain
        self.planning_prompt = PromptTemplate.from_template(
            """Create a security assessment plan based on:
            Company Research: {company_research}
            Latest Vulnerabilities: {latest_vulns}
            Previous Scan Results: {previous_results}
            
            Consider:
            1. Industry-specific threats
            2. Technology stack vulnerabilities
            3. Compliance requirements
            4. Past security incidents
            
            Create a prioritized testing plan that focuses on:
            1. Critical assets and crown jewels
            2. Known vulnerability patterns
            3. Compliance requirements
            4. Industry-specific attack vectors
            
            Format the response as a JSON object with these keys:
            - priority_scans: Ordered list of scans to run
            - scan_params: Specific parameters for each scan
            - success_criteria: What defines a successful scan
            - risk_thresholds: When to escalate findings
            """
        )
        self.planning_chain = self.planning_prompt | self.llm
        
    async def research_company(self, domain: str) -> Dict[str, Any]:
        """Research company information using web searches and analysis."""
        try:
            # Implement web search and data gathering here
            research = await self.research_chain.ainvoke({"domain": domain})
            return json.loads(research.content)
        except Exception as e:
            return {"error": str(e)}
            
    async def get_latest_vulnerabilities(self) -> List[Dict[str, Any]]:
        """Fetch latest vulnerability data from various sources."""
        try:
            # Example sources: NVD, CVE databases, security feeds
            async with aiohttp.ClientSession() as session:
                # Fetch from NVD
                nvd_url = "https://services.nvd.nist.gov/rest/json/cves/2.0"
                async with session.get(nvd_url) as response:
                    nvd_data = await response.json()
                
                # Process and return relevant vulnerabilities
                return self._process_vulnerability_data(nvd_data)
        except Exception as e:
            return [{"error": str(e)}]
            
    def _process_vulnerability_data(self, data: Dict) -> List[Dict[str, Any]]:
        """Process raw vulnerability data into useful format."""
        # Implementation here
        pass
        
    async def create_scan_plan(self, 
                             company_research: Dict[str, Any],
                             latest_vulns: List[Dict[str, Any]],
                             previous_results: List[Dict[str, Any]] = None) -> Dict[str, Any]:
        """Create an intelligent scan plan based on all available data."""
        try:
            plan = await self.planning_chain.ainvoke({
                "company_research": json.dumps(company_research),
                "latest_vulns": json.dumps(latest_vulns),
                "previous_results": json.dumps(previous_results if previous_results else [])
            })
            return json.loads(plan.content)
        except Exception as e:
            return {"error": str(e)}
            
    def create_report(self, 
                     scan_results: List[Dict[str, Any]], 
                     company_info: Dict[str, Any]) -> str:
        """Generate a comprehensive security report with visualizations."""
        try:
            # Create PowerPoint presentation
            prs = Presentation()
            
            # Add title slide
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_slide.shapes.title.text = f"Security Assessment Report: {company_info['company_info']['name']}"
            title_slide.placeholders[1].text = datetime.now().strftime("%Y-%m-%d")
            
            # Add executive summary
            summary_slide = prs.slides.add_slide(prs.slide_layouts[1])
            summary_slide.shapes.title.text = "Executive Summary"
            
            # Add findings visualization
            self._create_findings_chart(scan_results)
            findings_slide = prs.slides.add_slide(prs.slide_layouts[5])
            findings_slide.shapes.title.text = "Key Findings"
            findings_slide.shapes.add_picture("findings_chart.png", Inches(1), Inches(2))
            
            # Save presentation
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            report_path = f"data/scan_results/report_{timestamp}.pptx"
            prs.save(report_path)
            
            return report_path
            
        except Exception as e:
            return str(e)
            
    def _create_findings_chart(self, scan_results: List[Dict[str, Any]]):
        """Create visualization of scan findings."""
        # Implementation here using matplotlib/seaborn
        pass 